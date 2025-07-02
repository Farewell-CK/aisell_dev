import os
import json
import time
from typing import List, Dict, Optional, Union, Tuple # 添加 Optional, Union
from openai import OpenAI
from tqdm import tqdm
import base64
from pathlib import Path
import logging
import mimetypes # 用于更准确的MIME类型推断
import pandas as pd
from pptx import Presentation
from pptx.util import Inches # For creating dummy pptx
import PyPDF2
from docx import Document as DocxDocument # type: ignore # 重命名以避免与类名冲突
import subprocess
import tempfile
import shutil
from urllib.parse import urlparse # 用于解析URL
import requests # 用于下载URL内容
import io # 用于处理内存中的文件流

# 加载环境变量
from dotenv import load_dotenv
load_dotenv()

# --- 日志配置 ---
# 确保日志目录存在
log_dir = "logs"
if not os.path.exists(log_dir):
    os.makedirs(log_dir)
    logging.debug(f"Log directory created at {log_dir}.")
# 配置日志记录器
# 设置日志级别为INFO，意味着INFO、WARNING、ERROR、CRITICAL级别的日志都会被记录
# 日志格式包括时间戳、日志级别和日志消息
logging.basicConfig(
    level=logging.INFO, 
    format='%(asctime)s - %(levelname)s - %(module)s - %(funcName)s - %(lineno)d - %(message)s',
    handlers=[
        logging.FileHandler(os.path.join(log_dir, "summarizer.log"), encoding='utf-8'), # 输出到文件
        logging.StreamHandler() # 同时输出到控制台
    ]
)


class TextSummarizer:
    def __init__(self, api_key: str):
        """
        初始化 TextSummarizer。

        :param api_key: 用于访问文心一言 API 的密钥。
        """
        self.client = OpenAI(
            api_key=api_key,
            base_url="https://qianfan.baidubce.com/v2" # 百度文心API的兼容OpenAI接口地址
        )
        self.chunk_size = 4000  # 每个块的最大字符数
        self.overlap_size = 200  # 块之间的重叠字符数
        logging.info(f"TextSummarizer 初始化完成。Chunk size: {self.chunk_size}, Overlap size: {self.overlap_size}")

    def split_text(self, text: str) -> List[str]:
        """
        将文本分割成重叠的块。
        改进了start的更新逻辑，防止因分块过小和重叠过大导致的指针不前进或后退的问题。

        :param text: 待分割的原始文本。
        :return: 分割后的文本块列表。
        """
        chunks = []
        start = 0
        text_length = len(text)
        
        if text_length == 0:
            logging.warning("输入文本为空，无法分割。")
            return []

        logging.info(f"开始分割文本，总长度: {text_length} 字符。")
        while start < text_length:
            # 确定当前块的理论结束位置
            end = min(start + self.chunk_size, text_length)
            actual_chunk_end = end # 实际的块结束位置，可能会根据断句调整

            # 如果不是最后一个块 (即 end < text_length)，尝试在句子或段落边界处分割
            if end < text_length:
                # 在 [start, end) 区间内从后向前查找句号或换行符
                last_period = text.rfind('。', start, end)
                last_newline = text.rfind('\n', start, end)
                split_point = max(last_period, last_newline)
                
                # 如果找到了有效的分割点 (在 start 之后)
                if split_point > start:
                    actual_chunk_end = split_point + 1 # 分割点后一个字符作为块的结束
                    logging.debug(f"找到分割点 {split_point}，块 [{start}:{actual_chunk_end}]")
                else:
                    logging.debug(f"未在 [{start}:{end}] 找到合适的自然分割点，使用默认块结束位置 {actual_chunk_end}")
            
            chunks.append(text[start:actual_chunk_end])
            logging.debug(f"添加块: text[{start}:{actual_chunk_end}], 长度 {len(chunks[-1])}")

            # 如果当前块的结束位置已经是文本末尾，则结束分割
            if actual_chunk_end == text_length:
                logging.info("已到达文本末尾，分割结束。")
                break
            
            # 计算下一个块的起始位置，考虑重叠
            next_start_candidate = actual_chunk_end - self.overlap_size
            
            # 关键：确保 start 指针总是前进，防止无限循环
            if next_start_candidate <= start:
                logging.warning(
                    f"下一个起始点 ({next_start_candidate}) 小于等于当前起始点 ({start}). "
                    f"这可能因为块太短 (长度 {actual_chunk_end - start}) 而重叠 ({self.overlap_size}) 较大. "
                    f"将起始点移动到当前块的末尾 ({actual_chunk_end}) 以避免循环。"
                )
                start = actual_chunk_end # 强制前进，牺牲这一步的重叠
            else:
                start = next_start_candidate
            
            logging.debug(f"下一个块的起始位置: {start}")

        logging.info(f"文本分割完成，共 {len(chunks)} 个块。")
        return chunks

    def summarize_chunk(self, chunk: str, context: str = "") -> str:
        """
        使用文心一言API总结单个文本块。

        :param chunk: 待总结的文本块。
        :param context: 上下文信息，通常是前一个块的总结。
        :return: 总结后的文本，如果出错则返回空字符串。
        """
        prompt = f"""请对以下文本进行概述，提取关键信息，保持客观准确。要求：
1. 输出内容必须控制在50字左右
2. 只提取最核心的信息
3. 使用简洁的语言表达

                    上下文：{context}

                    待总结文本：
                    {chunk}

                    请提供简洁的概述："""

        try:
            logging.info(f"请求API总结，块长度: {len(chunk)}，上下文长度: {len(context)}")
            # logging.debug(f"发送给API的完整Prompt:\n{prompt}") # 如果需要详细调试，可以取消注释
            response = self.client.chat.completions.create(
                messages=[
                    {'role': 'system', 'content': '你是一个专业的文本总结助手，擅长提取文本的关键信息并生成简洁的概述。'},
                    {'role': 'user', 'content': prompt}
                ],
                model="ernie-4.5-turbo-128k", # 使用文心3.5 8k模型
                temperature=0.3,      #较低的temperature使输出更稳定和集中
                top_p=0.8             # top_p核采样
            )
            summary = response.choices[0].message.content
            logging.info(f"API调用成功，返回总结长度: {len(summary)}")
            # logging.debug(f"API返回的总结内容: {summary[:200]}...") # 打印部分返回内容
            return summary
        except Exception as e:
            logging.error(f"API调用出错: {str(e)}", exc_info=True) # exc_info=True会记录堆栈跟踪
            time.sleep(5)  # 出错后等待5秒再尝试（如果是在循环中）
            return "" # 返回空字符串表示此块总结失败

    def _hierarchical_summarize(self, text_to_summarize: str, base_prompt_context: str, level: int = 0) -> str:
        """
        分层总结文本。如果文本过长，则分块总结，然后总结这些摘要。

        :param text_to_summarize: 需要总结的文本。
        :param base_prompt_context: 基础的提示词上下文。
        :param level: 当前的总结层级（用于日志）。
        :return: 最终的总结文本。
        """
        logging.info(f"进入分层总结 (层级 {level})，待总结文本长度: {len(text_to_summarize)}")

        # 如果文本长度小于等于块大小的某个比例（例如1.2倍，为提示词本身留足空间），则直接总结
        # 这里的 self.chunk_size 是针对原始文本分割的，对于总结的总结，可以稍微放宽，但仍需注意API限制
        # 实际上，API的token限制才是硬限制。4000字符约1000-3000 token。
        # 一个更安全的做法是检查字符数是否在一个安全的范围内，比如略小于 self.chunk_size
        if len(text_to_summarize) <= self.chunk_size * 1.1: # 乘以1.1是给prompt中的其他文字留余地
            logging.info(f"文本长度适中 (层级 {level})，直接进行单次总结。")
            return self.summarize_chunk(text_to_summarize, base_prompt_context)
        else:
            logging.info(f"文本过长 (层级 {level}，长度 {len(text_to_summarize)})，进行分块总结。")
            chunks = self.split_text(text_to_summarize) # 使用同样的分割逻辑
            if not chunks:
                logging.warning(f"分层总结 (层级 {level}) 时未能将文本分割成块。")
                return "未能生成有效的总结（分块失败）。"

            summaries = []
            current_context_for_summary_chunks = "" # 用于总结摘要块的上下文

            logging.info(f"文本 (层级 {level}) 被分割为 {len(chunks)} 个子块进行总结。")
            for i, chunk in enumerate(tqdm(chunks, desc=f"处理层级 {level} 的总结块")):
                logging.info(f"正在总结子块 {i+1}/{len(chunks)} (层级 {level})")
                # 为摘要的摘要构建提示词
                chunk_prompt = f"{base_prompt_context} 这是其中一部分，请总结:"
                summary = self.summarize_chunk(chunk, current_context_for_summary_chunks + "\n" + chunk_prompt if current_context_for_summary_chunks else chunk_prompt)
                if summary:
                    summaries.append(summary)
                    current_context_for_summary_chunks = summary # 使用当前子摘要作为下一个子摘要的上下文
                else:
                    logging.warning(f"子块 {i+1}/{len(chunks)} (层级 {level}) 未能生成总结。")
                time.sleep(1) # 控制API调用频率
            
            if not summaries:
                logging.warning(f"分层总结 (层级 {level}) 未能生成任何子摘要。")
                return "未能生成有效的总结（子摘要为空）。"
            
            combined_summaries = "\n\n".join(summaries)
            logging.info(f"分层总结 (层级 {level}) 后的合并摘要长度: {len(combined_summaries)}")
            
            # 递归调用，对合并后的摘要进行再总结
            return self._hierarchical_summarize(combined_summaries, f"这是对上一级总结内容的再总结 (层级 {level+1})，请进一步精炼：", level + 1)


    def process_file(self, file_path: str):
        """
        处理整个文件并生成总结。

        :param file_path: 输入文件的路径。
        :param output_path: 输出总结文件的路径。
        """
        logging.info(f"开始处理文件: {file_path}")
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            logging.info(f"文件读取成功，原始文本长度: {len(text)} 字符")
        except Exception as e:
            logging.error(f"读取文件 '{file_path}' 出错: {str(e)}", exc_info=True)
            return

        if not text.strip():
            logging.warning(f"文件 '{file_path}' 内容为空或只包含空白字符。")
            return "输入文件为空，无法生成总结。"
            '''
            try:
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write("输入文件为空，无法生成总结。")
                logging.info(f"空文件提示已保存到: {output_path}")
            except Exception as e:
                logging.error(f"保存空文件提示到 '{output_path}' 出错: {str(e)}", exc_info=True)
            return
            '''
        initial_chunks = self.split_text(text)

        if not initial_chunks:
            logging.warning("未能将原始文本分割成任何块。处理中止。")
            return "未能将原始文本分割成任何块。处理中止。"

        '''
        if not initial_chunks:
            logging.warning("未能将原始文本分割成任何块。处理中止。")
            try:
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write("文本分割失败，无法生成总结。")
            except Exception as e:
                logging.error(f"保存分割失败提示到 '{output_path}' 出错: {str(e)}", exc_info=True)
            return
        '''   

        logging.info(f"原始文本已分割为 {len(initial_chunks)} 个块。")
        
        intermediate_summaries = []
        current_chunk_context = "" # 用于指导下一个块总结的上下文

        for i, chunk in enumerate(tqdm(initial_chunks, desc="处理初始文本块")):
            logging.info(f"正在处理初始块 {i+1}/{len(initial_chunks)}，长度 {len(chunk)}")
            summary = self.summarize_chunk(chunk, current_chunk_context)
            if summary:
                intermediate_summaries.append(summary)
                current_chunk_context = summary  # 使用当前块的总结作为下一个块的上下文
                logging.info(f"初始块 {i+1} 总结完成。摘要长度: {len(summary)}。上下文已更新。")
            else:
                logging.warning(f"初始块 {i+1} 未能生成总结。")
            time.sleep(1)  # 控制API调用频率

        if not intermediate_summaries:
            logging.warning("未能生成任何初始文本块的总结，无法进行最终总结。")
            return "未能对任何文本块生成初步总结。"
            '''
            try:
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write("未能对任何文本块生成初步总结。")
                logging.info(f"无初步总结提示已保存到: {output_path}")
            except Exception as e:
                logging.error(f"保存无初步总结提示到 '{output_path}' 出错: {str(e)}", exc_info=True)
            return
            '''

        logging.info("所有初始块总结完成，正在合并初步总结内容...")
        combined_initial_summaries = "\n\n".join(intermediate_summaries)
        logging.info(f"合并后的初步总结文本长度: {len(combined_initial_summaries)} 字符")

        # 对合并后的初步总结进行分层总结
        logging.info("开始最终的分层总结流程...")
        final_summary_prompt = "这是对所有文本片段初步总结的整合，请基于这些内容生成一个全面且精炼的最终概述(字数在100字以内,是对文件内容的描述)："
        final_summary = self._hierarchical_summarize(combined_initial_summaries, final_summary_prompt)
        return final_summary
        # 保存结果
        '''
        try:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(final_summary)
            logging.info(f"最终总结已生成并保存到: {output_path}")
            logging.info(f"最终总结内容 (前500字符): {final_summary[:500]}...")
        except Exception as e:
            logging.error(f"保存最终总结到 '{output_path}' 出错: {str(e)}", exc_info=True)
        '''

class ImageSummarizer:
    def __init__(self, api_key: str):
        """
        初始化 ImageSummarizer。
        :param api_key: 用于访问 Ernie API 的密钥。
        """
        self.client = OpenAI(
            api_key=api_key,
            base_url="https://qianfan.baidubce.com/v2"
        )
        self.model_name = "ernie-4.5-turbo-vl-preview" # 您指定的模型
        logging.info(f"ImageSummarizer 初始化完成。API base_url: {self.client.base_url}, Model: {self.model_name}")

    def encode_image(self, image_path: str) -> Optional[str]:
        """
        将本地图片文件转换为base64编码的字符串。
        :param image_path: 本地图片文件的路径。
        :return: Base64编码的字符串，如果出错则返回None。
        """
        logging.debug(f"开始对图片进行Base64编码: {image_path}")
        try:
            with open(image_path, "rb") as image_file:
                encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
            logging.info(f"图片 '{image_path}' Base64编码成功。编码后数据长度(部分): {len(encoded_string[:100])}...")
            return encoded_string
        except FileNotFoundError:
            logging.error(f"图片文件未找到: {image_path}")
            return None
        except Exception as e:
            logging.error(f"对图片 '{image_path}'进行Base64编码时出错: {e}", exc_info=True)
            return None

    def get_image_url_dict(self, image_path_or_url: str) -> Optional[Dict[str, str]]:
        """
        根据输入是URL还是本地路径，生成API所需的image_url字典。
        对于本地文件，会进行Base64编码并生成Data URL。
        :param image_path_or_url: 图片的URL或本地文件路径。
        :return: API兼容的image_url字典，如果本地文件处理失败则返回None。
        """
        logging.debug(f"正在准备图片URL字典: {image_path_or_url}")
        if image_path_or_url.startswith(('http://', 'https://')):
            logging.info(f"输入为URL，直接使用: {image_path_or_url}")
            return {"url": image_path_or_url}
        else:
            # 本地文件路径，进行Base64编码
            logging.info(f"输入为本地文件路径，将进行Base64编码: {image_path_or_url}")
            if not os.path.exists(image_path_or_url):
                logging.error(f"本地图片文件不存在: {image_path_or_url}")
                return None
            
            base64_image = self.encode_image(image_path_or_url)
            if base64_image is None:
                return None # encode_image内部已记录错误

            # 推断MIME类型
            mime_type, _ = mimetypes.guess_type(image_path_or_url)
            if not mime_type:
                mime_type = 'image/jpeg' # 如果无法推断，默认为jpeg（与原逻辑保持一致，但可优化）
                logging.warning(f"无法准确推断图片 '{image_path_or_url}' 的MIME类型，将默认使用 'image/jpeg'。")
            else:
                logging.info(f"推断出图片 '{image_path_or_url}' 的MIME类型为: {mime_type}")
            
            data_url = f"data:{mime_type};base64,{base64_image}"
            return {"url": data_url}

    def summarize_single_image(self, image_path: str, prompt: Optional[str] = None) -> str:
        """
        使用 LLM 总结单张图片的内容。
        :param image_path: 图片的URL或本地文件路径。
        :param prompt: 用户自定义的总结提示词。
        :return: LLM 生成的图片总结文本，或错误信息。
        """
        logging.info(f"开始处理单张图片总结: {image_path}")

        image_url_data = self.get_image_url_dict(image_path)
        if image_url_data is None:
            error_msg = f"无法获取或处理图片 '{image_path}'。"
            logging.error(error_msg)
            return error_msg

        if prompt is None:
            prompt = "请用中文详细描述这张图片的内容，包括但不限于：画面中的主要物体、人物（如有，描述其姿态、表情、衣着）、场景环境、光线氛围、可能的事件或故事背景，以及图片的整体风格和给人的感受。（50字以内）"
            logging.debug("使用默认提示词进行单图总结。")
        else:
            logging.debug(f"使用自定义提示词进行单图总结: '{prompt[:100]}...'")
            
        try:
            logging.info(f"向模型 '{self.model_name}' 发送单图总结请求。图片来源: {image_path}")
            messages_payload = [
                {
                    "role": "system",
                    "content": "你是一位专业的图像内容分析师，能够洞察图片的每一个细节并用生动的语言描述出来。(100字以内)"
                },
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": image_url_data}
                    ]
                }
            ]
            
            response = self.client.chat.completions.create(
                model=self.model_name,
                messages=messages_payload, # type: ignore
                temperature=0.3,
                top_p=0.8
            )
            
            summary = response.choices[0].message.content
            logging.info(f"成功从模型获取单图总结。总结长度: {len(summary) if summary else 0} 字符。")
            return summary if summary else "未能从模型获取有效的总结内容。"
            
        except Exception as e:
            logging.error(f"使用模型分析单张图片 '{image_path}' 时出错: {e}", exc_info=True)
            return f"分析图片时发生错误: {str(e)}"

    def compare_images(self, image_paths: List[str], prompt: Optional[str] = None) -> str:
        """
        使用 LLM 比较多张图片的异同。
        :param image_paths: 包含多个图片URL或本地文件路径的列表。
        :param prompt: 用户自定义的比较提示词。
        :return: LLM 生成的图片比较结果文本，或错误信息。
        """
        logging.info(f"开始处理图片比较，共 {len(image_paths)} 张图片。")
        if len(image_paths) < 2:
            logging.warning("需要至少两张图片进行比较，当前提供数量不足。")
            return "错误：需要至少两张图片进行比较。"
            
        if prompt is None:
            prompt = "请用中文详细对比分析这些图片之间的异同点，关注以下方面：1. 主题与内容元素；2. 构图与视角；3. 色彩与光影；4. 风格与氛围；5. 可能的关联或序列关系。"
            logging.debug("使用默认提示词进行多图比较。")
        else:
            logging.debug(f"使用自定义提示词进行多图比较: '{prompt[:100]}...'")
            
        user_content_list: List[Dict[str, Union[str, Dict[str, str]]]] = [{"type": "text", "text": prompt}]
        valid_images_for_api = 0

        for image_path in image_paths:
            logging.debug(f"为比较准备图片: {image_path}")
            image_url_data = self.get_image_url_dict(image_path)
            if image_url_data:
                user_content_list.append({"type": "image_url", "image_url": image_url_data})
                valid_images_for_api += 1
            else:
                logging.warning(f"无法处理图片 '{image_path}'，将从比较中排除。")
        
        if valid_images_for_api < 2:
            logging.error(f"有效图片数量 ({valid_images_for_api}) 不足两个，无法进行比较。")
            return f"错误：未能准备足够的有效图片进行比较 (至少需要2张，实际有效 {valid_images_for_api} 张)。"

        try:
            logging.info(f"向模型 '{self.model_name}' 发送 {valid_images_for_api} 张图片的多图比较请求。")
            messages_payload = [
                {
                    "role": "system",
                    "content": "你是一位顶级的图像对比分析专家，能够洞察多张图片间的细微差别和深层联系。"
                },
                {
                    "role": "user",
                    "content": user_content_list # type: ignore
                }
            ]
            response = self.client.chat.completions.create(
                model=self.model_name,
                messages=messages_payload, # type: ignore
                temperature=0.3,
                top_p=0.8
            )
            comparison_summary = response.choices[0].message.content
            logging.info(f"成功从模型获取多图比较结果。结果长度: {len(comparison_summary) if comparison_summary else 0} 字符。")
            return comparison_summary if comparison_summary else "未能从模型获取有效的比较结果。"
        except Exception as e:
            logging.error(f"使用模型比较多张图片时出错: {e}", exc_info=True)
            return f"比较图片时发生错误: {str(e)}"

    def process_image_directory(self, directory_path: str, output_file: str):
        """
        处理目录中的所有图片文件，对每张图片生成总结，并将所有总结写入一个报告文件。
        :param directory_path: 包含图片文件的目录路径。
        :param output_file: 保存所有总结报告的文件路径。
        """
        logging.info(f"开始批量处理图片目录: {directory_path}，输出到: {output_file}")
        try:
            dir_path = Path(directory_path)
            if not dir_path.is_dir():
                logging.error(f"提供的路径 '{directory_path}' 不是一个有效的目录。")
                return

            image_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff'} # 扩展了支持的格式
            image_files = [
                f for f in dir_path.rglob('*') # 使用 rglob 递归查找
                if f.is_file() and f.suffix.lower() in image_extensions
            ]
            
            if not image_files:
                logging.warning(f"在目录 '{directory_path}' 及其子目录中没有找到支持的图片文件。")
                Path(output_file).write_text(f"在目录 {directory_path} 中没有找到支持的图片文件。\n", encoding='utf-8')
                return
            
            logging.info(f"在 '{directory_path}' 中找到 {len(image_files)} 个图片文件: {[str(f) for f in image_files]}")

            all_summaries_content: List[str] = []
            successful_summaries = 0
            
            for image_file_path in tqdm(image_files, desc="批量处理目录中的图片"):
                logging.info(f"正在处理图片文件: {str(image_file_path)}")
                summary = self.summarize_single_image(str(image_file_path))
                
                if summary and not summary.startswith("错误：") and not summary.startswith("无法获取或处理图片") and not summary.startswith("分析图片时发生错误"):
                    all_summaries_content.append(f"--- 图片: {image_file_path.name} ---\n{summary}\n\n")
                    successful_summaries +=1
                else:
                    all_summaries_content.append(f"--- 图片: {image_file_path.name} ---\n总结失败或文件无效。\n详细信息: {summary}\n\n")
                
                if len(image_files) > 1 and image_file_path != image_files[-1]: # 如果不止一个文件且不是最后一个，则等待
                    logging.debug("等待1秒以控制API调用频率...")
                    time.sleep(1)

            report_header = f"图片批量处理总结报告\n处理图片总数: {len(image_files)}\n成功总结数量: {successful_summaries}\n\n"
            final_report_content = report_header + "".join(all_summaries_content)

            Path(output_file).write_text(final_report_content, encoding='utf-8')
            logging.info(f"所有图片的总结报告已保存到: {output_file}")

        except Exception as e:
            logging.error(f"处理图片目录 '{directory_path}' 时发生意外错误: {e}", exc_info=True)

class TableSummarizer:
    def __init__(self, api_key: str):
        """
        初始化 TableSummarizer。

        :param api_key: 用于访问文心一言 API 的密钥。
        """
        self.client = OpenAI(
            api_key=api_key,
            base_url="https://qianfan.baidubce.com/v2"
        )
        logging.info(f"TableSummarizer 初始化完成。API base_url: https://qianfan.baidubce.com/v2")

    def is_valid_url(self, url: str) -> bool:
        """简单的URL格式检查"""
        try:
            result = urlparse(url)
            return all([result.scheme in ['http', 'https'], result.netloc])
        except:
            return False

    def read_table(self, file_path_or_url: str) -> pd.DataFrame:
        """
        读取表格文件 (xls, xlsx, csv)，支持本地路径和URL。

        :param file_path_or_url: 表格文件的本地路径或URL。
        :return: 读取到的 pandas DataFrame，如果出错则返回空 DataFrame。
        """
        logging.info(f"尝试读取表格源: {file_path_or_url}")

        if self.is_valid_url(file_path_or_url):
            logging.info(f"检测到输入为 URL: {file_path_or_url}")
            try:
                # 从URL路径中提取文件名和扩展名以猜测类型
                parsed_url = urlparse(file_path_or_url)
                path_part = parsed_url.path
                filename = os.path.basename(path_part)
                file_ext = os.path.splitext(filename)[1].lower()
                
                logging.info(f"从URL推断文件名: '{filename}', 扩展名: '{file_ext}'")

                headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537.3'}
                response = requests.get(file_path_or_url, headers=headers, stream=True, timeout=30) # stream=True for large files, timeout
                response.raise_for_status()  # 如果HTTP请求返回不成功的状态码，则抛出HTTPError异常
                
                logging.info(f"URL 内容下载成功 (状态码: {response.status_code})。尝试解析...")

                if file_ext in ['.xls', '.xlsx']:
                    # Excel文件需要二进制内容
                    content_bytes = io.BytesIO(response.content)
                    df = pd.read_excel(content_bytes)
                    logging.info(f"成功从 URL 读取 Excel 内容: {file_path_or_url}")
                    return df
                elif file_ext == '.csv':
                    # CSV文件可以是文本，但用BytesIO然后decode更稳妥处理编码问题
                    # 或者尝试直接用 response.text，但要注意编码
                    try:
                        # 尝试使用 utf-8 解码
                        content_text = response.content.decode('utf-8')
                        df = pd.read_csv(io.StringIO(content_text))
                        logging.info(f"成功从 URL 读取 CSV 内容 (UTF-8): {file_path_or_url}")
                    except UnicodeDecodeError:
                        logging.warning(f"URL '{file_path_or_url}' 的 CSV 内容 UTF-8 解码失败，尝试 GBK 解码...")
                        try:
                            content_text_gbk = response.content.decode('gbk')
                            df = pd.read_csv(io.StringIO(content_text_gbk))
                            logging.info(f"成功从 URL 读取 CSV 内容 (GBK): {file_path_or_url}")
                        except Exception as e_gbk:
                            logging.error(f"URL '{file_path_or_url}' 的 CSV 内容 GBK 解码也失败: {e_gbk}. 尝试直接用BytesIO.")
                            # 如果解码困难，pandas可以直接从BytesIO读取csv，它会尝试推断编码
                            content_bytes = io.BytesIO(response.content)
                            df = pd.read_csv(content_bytes)
                            logging.info(f"成功从 URL 读取 CSV 内容 (Pandas自动推断编码): {file_path_or_url}")
                    return df
                else:
                    logging.warning(f"无法从URL的文件扩展名 '{file_ext}' 确定支持的表格类型。URL: {file_path_or_url}")
                    # 也可以尝试从 Content-Type header 推断，但这里简化处理
                    return pd.DataFrame()

            except requests.exceptions.RequestException as e_req:
                logging.error(f"下载URL '{file_path_or_url}' 时发生网络错误: {e_req}", exc_info=True)
                return pd.DataFrame()
            except Exception as e:
                logging.error(f"处理URL '{file_path_or_url}' 内容时出错: {e}", exc_info=True)
                return pd.DataFrame()
        
        else: # 处理本地文件路径
            logging.info(f"输入为本地文件路径: {file_path_or_url}")
            file_ext = os.path.splitext(file_path_or_url)[1].lower()
            logging.info(f"本地文件格式: {file_ext}")
            try:
                if file_ext in ['.xls', '.xlsx']:
                    df = pd.read_excel(file_path_or_url)
                    logging.info(f"成功读取本地 Excel 文件: {file_path_or_url}")
                    return df
                elif file_ext == '.csv':
                    df = pd.read_csv(file_path_or_url)
                    logging.info(f"成功读取本地 CSV 文件: {file_path_or_url}")
                    return df
                else:
                    logging.warning(f"不支持的本地文件格式: {file_ext} 文件路径: {file_path_or_url}")
                    return pd.DataFrame()
            except Exception as e:
                logging.error(f"读取本地文件 '{file_path_or_url}' 出错: {str(e)}", exc_info=True)
                return pd.DataFrame()
            
    def get_table_info(self, df: pd.DataFrame) -> Dict:
        """
        获取表格 DataFrame 的基本信息和统计数据。

        :param df: 输入的 pandas DataFrame。
        :return: 包含表格信息的字典。
        """
        if df.empty:
            logging.warning("输入 DataFrame 为空，无法提取表格信息。")
            return {
                "错误": "输入DataFrame为空"
            }

        logging.info(f"开始提取表格信息，行数: {len(df)}, 列数: {len(df.columns)}")
        info = {
            "行数": len(df),
            "列数": len(df.columns),
            "列名": list(df.columns),
            "数据类型": {col: str(df[col].dtype) for col in df.columns},
            "空值总数": int(df.isnull().sum().sum()), # 表格总空值数
            "每列空值数": {col: int(df[col].isnull().sum()) for col in df.columns},
            "数值列统计": {}
        }

        numeric_cols = df.select_dtypes(include=['int64', 'float64', 'int32', 'float32']).columns
        logging.info(f"找到数值列: {list(numeric_cols)}")
        for col in numeric_cols:
            if df[col].isnull().all(): # 如果整列都是NaN
                logging.warning(f"数值列 '{col}' 完全为空，统计值将为 NaN 或引发错误。")
                info["数值列统计"][col] = {
                    "平均值": float('nan'), "中位数": float('nan'),
                    "最大值": float('nan'), "最小值": float('nan'),
                    "标准差": float('nan'), "备注": "该列完全为空值"
                }
            else:
                info["数值列统计"][col] = {
                    "平均值": float(df[col].mean()),
                    "中位数": float(df[col].median()),
                    "最大值": float(df[col].max()),
                    "最小值": float(df[col].min()),
                    "标准差": float(df[col].std())
                }
        logging.info("表格信息提取完成。")
        return info

    def summarize_table(self, file_path: str, custom_prompt: str = None) -> str:
        """
        使用 LLM 总结单个表格的内容。

        :param file_path: 表格文件的路径。
        :param custom_prompt: 用户自定义的分析提示词。
        :return: LLM 生成的表格总结文本。
        """
        logging.info(f"开始总结表格: {file_path}")
        df = self.read_table(file_path)
        if df.empty:
            logging.warning(f"无法读取或表格为空: {file_path}，总结中止。")
            return "无法读取表格文件或表格内容为空。"

        table_info = self.get_table_info(df)
        if "错误" in table_info: # 来自 get_table_info 的错误
             return f"提取表格信息失败: {table_info['错误']}"

        try:
            # 将表格信息和数据示例转换为字符串
            # 注意：对于非常大的表格信息或宽表格，这里可能产生很长的字符串
            info_str = json.dumps(table_info, ensure_ascii=False, indent=2, default=str) # default=str 处理 NaN 等
            sample_data_df = df.head() # 取前5行
            sample_data_str = sample_data_df.to_string()
            logging.info(f"表格信息字符串长度: {len(info_str)}")
            logging.info(f"数据示例字符串长度: {len(sample_data_str)}")
            # 考虑截断过长的字符串
            MAX_INFO_LEN = 3000 # 示例最大长度，根据模型token限制调整
            MAX_SAMPLE_LEN = 2000 # 示例最大长度
            if len(info_str) > MAX_INFO_LEN:
                logging.warning(f"表格信息过长，已截断至 {MAX_INFO_LEN} 字符。")
                info_str = info_str[:MAX_INFO_LEN] + "\n... [信息已截断]"
            if len(sample_data_str) > MAX_SAMPLE_LEN:
                logging.warning(f"数据示例过长，已截断至 {MAX_SAMPLE_LEN} 字符。")
                sample_data_str = sample_data_str[:MAX_SAMPLE_LEN] + "\n... [数据已截断]"

        except Exception as e:
            logging.error(f"序列化表格信息或数据示例时出错: {str(e)}", exc_info=True)
            return "处理表格信息时发生内部错误。"


        prompt_to_use = custom_prompt if custom_prompt else """请基于以下提供的表格元信息和数据抽样，用中文分析这个表格的内容，包括：
1.  表格的基本结构和关键列描述。
2.  数据的主要特征、分布和潜在趋势（例如，数值列的集中趋势、离散程度，分类列的众数等）。
3.  指出现有数据中可能存在的数据质量问题（例如，大量的空值、异常值、不一致的数据格式等）。
4.  根据现有信息，总结出一些关键的洞察或发现。
5.  如果数据适合某种类型的分析（如时间序列、分类、回归等），请指出。

要求：
1. 输出内容必须控制在50字左右
2. 只提取最核心的信息
3. 使用简洁的语言表达
"""

        user_content = f"{prompt_to_use}\n\n# 表格元信息:\n```json\n{info_str}\n```\n\n# 数据抽样 (前 {len(sample_data_df)} 行):\n```\n{sample_data_str}\n```"
        logging.debug(f"发送给 LLM 的 User Content (部分): {user_content[:500]}...")

        try:
            logging.info("向 LLM 发送请求进行表格总结...")
            response = self.client.chat.completions.create(
                model="ernie-4.5-8k-preview", 
                messages=[
                    {"role": "system", "content": "你是一个专业的数据分析师，擅长从表格的结构信息和数据样本中提炼洞察。"},
                    {"role": "user", "content": user_content}
                ],
                temperature=0.3,
                top_p=0.8
            )
            summary = response.choices[0].message.content
            logging.info(f"LLM 表格总结接收成功，长度: {len(summary)}")
            return summary
        except Exception as e:
            logging.error(f"调用 LLM 分析表格时出错: {str(e)}", exc_info=True)
            return "分析表格时与AI服务通信失败。"

    def compare_tables(self, file_paths: List[str], custom_prompt: str = None) -> str:
        """
        使用 LLM 比较多个表格的异同。

        :param file_paths: 包含多个表格文件路径的列表。
        :param custom_prompt: 用户自定义的比较提示词。
        :return: LLM 生成的表格比较结果文本。
        """
        logging.info(f"开始比较表格: {file_paths}")
        if len(file_paths) < 2:
            logging.warning("需要至少两个表格文件进行比较。")
            return "错误：需要至少两个表格文件进行比较。"

        tables_metadata = []
        for file_path in file_paths:
            df = self.read_table(file_path)
            if not df.empty:
                info = self.get_table_info(df)
                if "错误" not in info:
                    info['文件名'] = os.path.basename(file_path)
                    tables_metadata.append(info)
                else:
                    logging.warning(f"跳过文件 {file_path}，因提取信息失败。")
            else:
                logging.warning(f"跳过文件 {file_path}，因读取失败或为空。")


        if not tables_metadata:
            logging.error("无法读取任何有效的表格文件进行比较。")
            return "错误：无法读取任何有效的表格文件。"
        if len(tables_metadata) < 2:
            logging.warning(f"只成功读取了 {len(tables_metadata)} 个表格，无法进行有效比较。")
            return f"错误：只成功读取了 {len(tables_metadata)} 个有效表格，无法比较。"


        try:
            metadata_str = json.dumps(tables_metadata, ensure_ascii=False, indent=2, default=str)
            logging.info(f"多表格元信息字符串长度: {len(metadata_str)}")
            MAX_META_LEN = 7000 # 多个表格信息可能非常长，根据模型调整
            if len(metadata_str) > MAX_META_LEN:
                logging.warning(f"多表格元信息过长，已截断至 {MAX_META_LEN} 字符。")
                metadata_str = metadata_str[:MAX_META_LEN] + "\n... [元信息已截断]"
        except Exception as e:
            logging.error(f"序列化多表格元信息时出错: {str(e)}", exc_info=True)
            return "处理多表格元信息时发生内部错误。"

        prompt_to_use = custom_prompt if custom_prompt else """请基于以下提供的多个表格的元信息，用中文比较这些表格的异同点，分析内容应包括但不限于：
                1.  **结构对比**：例如列名、列数量、数据类型的相似与差异。
                2.  **数据规模对比**：例如行数的对比。
                3.  **数据特征对比**：根据数值列统计（均值、中位数、极差等）和空值情况，对比它们的数据分布和数据质量的异同。
                4.  **潜在关联或差异总结**：基于以上对比，总结这些表格可能反映的共同模式、显著差异或潜在的数据问题。"""

        user_content = f"{prompt_to_use}\n\n# 各表格元信息汇总:\n```json\n{metadata_str}\n```"
        logging.debug(f"发送给 LLM 的 User Content (部分): {user_content[:500]}...")

        try:
            logging.info("向 LLM 发送请求进行表格比较...")
            response = self.client.chat.completions.create(
                model="ernie-4.0-8k-preview",
                messages=[
                    {"role": "system", "content": "你是一个专业的数据对比分析师，擅长从多个表格的结构和统计信息中找出异同和关联。"},
                    {"role": "user", "content": user_content}
                ],
                temperature=0.3,
                top_p=0.8
            )
            comparison_summary = response.choices[0].message.content
            logging.info(f"LLM 表格比较结果接收成功，长度: {len(comparison_summary)}")
            return comparison_summary
        except Exception as e:
            logging.error(f"调用 LLM 比较表格时出错: {str(e)}", exc_info=True)
            return "比较表格时与AI服务通信失败。"

    def process_directory(self, directory_path: str, output_file: str):
        """
        处理目录中的所有表格文件 (xls, xlsx, csv)，并对每个文件生成总结。

        :param directory_path: 包含表格文件的目录路径。
        :param output_file: 保存所有总结报告的文件路径。
        """
        logging.info(f"开始处理目录: {directory_path}，输出到: {output_file}")
        try:
            dir_path = Path(directory_path)
            if not dir_path.is_dir():
                logging.error(f"提供的路径 '{directory_path}' 不是一个有效的目录。")
                return

            table_files = []
            # 使用 rglob 进行递归搜索
            for ext in ['.xls', '.xlsx', '.csv']:
                table_files.extend(list(dir_path.rglob(f'*{ext}'))) # 使用 rglob
            
            if not table_files:
                logging.warning(f"在目录 '{directory_path}' 及其子目录中没有找到支持的表格文件。")
                return
            
            logging.info(f"在 '{directory_path}' 中找到 {len(table_files)} 个表格文件: {[str(f) for f in table_files]}")

            all_summaries_content = []
            for table_file_path in tqdm(table_files, desc="逐个处理目录中的表格"):
                logging.info(f"正在处理文件: {str(table_file_path)}")
                summary = self.summarize_table(str(table_file_path))
                if summary and not summary.startswith("无法读取") and not summary.startswith("错误："):
                    all_summaries_content.append(f"--- 文件: {table_file_path.name} ---\n{summary}\n\n")
                else:
                    all_summaries_content.append(f"--- 文件: {table_file_path.name} ---\n总结失败或文件无效。\n详细原因: {summary}\n\n")
                time.sleep(1)  # 控制API调用频率，避免过于频繁

            if not all_summaries_content:
                logging.info("未能对目录中的任何文件生成有效总结。")
                report_content = "未能对目录中的任何文件生成有效总结。"
            else:
                report_content = "".join(all_summaries_content)

            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(report_content)
            logging.info(f"所有表格的总结报告已保存到: {output_file}")

        except Exception as e:
            logging.error(f"处理目录 '{directory_path}' 时发生意外错误: {str(e)}", exc_info=True)

# --- 辅助函数：创建示例文件 ---
def create_dummy_excel(file_path, num_rows=10):
    if not os.path.exists(file_path):
        logging.info(f"创建示例 Excel 文件: {file_path}")
        data = {
            'ID': range(1, num_rows + 1),
            '名称': [f'项目_{i}' for i in range(1, num_rows + 1)],
            '数值A': [i * 10.5 for i in range(1, num_rows + 1)],
            '数值B': [i * 100 if i % 2 == 0 else pd.NA for i in range(1, num_rows + 1)], # 包含空值
            '类别': ['A' if i % 3 == 0 else 'B' if i % 3 == 1 else 'C' for i in range(1, num_rows + 1)]
        }
        df = pd.DataFrame(data)
        # 确保目录存在
        os.makedirs(os.path.dirname(file_path) or '.', exist_ok=True)
        df.to_excel(file_path, index=False)


class PPTSummarizer:
    def __init__(self, api_key: str):
        """
        初始化 PPTSummarizer。
        :param api_key: 用于访问文心一言 API 的密钥。
        """
        self.client = OpenAI(
            api_key=api_key,
            base_url="https://qianfan.baidubce.com/v2"
        )
        self.model_name = "ernie-4.5-8k-preview" # 您指定的模型
        logging.info(f"PPTSummarizer 初始化完成。API base_url: {self.client.base_url}, Model: {self.model_name}")
        logging.info("确保 LibreOffice (soffice) 已安装并在系统 PATH 中，以便支持 .ppt 文件转换。")

    def _is_valid_url(self, url: str) -> bool:
        """简单的URL格式检查，判断是否为 HTTP/HTTPS URL"""
        try:
            result = urlparse(url)
            return all([result.scheme in ['http', 'https'], result.netloc])
        except Exception:
            return False

    def convert_ppt_to_pptx(self, ppt_path: str) -> Tuple[Optional[str], Optional[str]]:
        """
        使用LibreOffice将PPT转换为PPTX。
        返回 (converted_file_path, temp_dir_path) 或 (None, None) 如果失败。
        """
        logging.info(f"开始将PPT文件 '{ppt_path}' 转换为PPTX格式。")
        temp_dir_for_conversion = None # 用于存放 soffice 转换结果的临时目录
        try:
            # 这个临时目录是 soffice 输出转换后 .pptx 文件的地方
            temp_dir_for_conversion = tempfile.mkdtemp()
            logging.debug(f"创建转换输出临时目录: {temp_dir_for_conversion}")

            cmd = [
                'soffice',
                '--headless',
                '--invisible',
                '--convert-to', 'pptx',
                '--outdir', temp_dir_for_conversion, # soffice 输出到此目录
                ppt_path # 输入的 .ppt 文件路径 (可以是下载到本地的临时.ppt文件)
            ]
            logging.info(f"执行转换命令: {' '.join(cmd)}")
            process = subprocess.run(cmd, check=True, capture_output=True, timeout=180) # 增加超时到3分钟
            logging.info(f"LibreOffice转换成功。Stdout: {process.stdout.decode('utf-8', 'ignore')}")
            if process.stderr:
                 logging.warning(f"LibreOffice转换 stderr: {process.stderr.decode('utf-8', 'ignore')}")

            original_basename = os.path.basename(ppt_path) # 基于输入ppt文件名确定输出名
            converted_filename = os.path.splitext(original_basename)[0] + '.pptx'
            converted_file_path = os.path.join(temp_dir_for_conversion, converted_filename)

            if os.path.exists(converted_file_path):
                logging.info(f"PPTX文件已成功转换并保存到: {converted_file_path}")
                return converted_file_path, temp_dir_for_conversion
            else:
                logging.error(f"转换后的PPTX文件 '{converted_file_path}' 未找到。输入路径: {ppt_path}")
                # 清理为 soffice 输出创建的临时目录
                shutil.rmtree(temp_dir_for_conversion)
                return None, None
        except FileNotFoundError:
            logging.error("LibreOffice (soffice) 未找到。请确保已安装并配置在系统PATH中。", exc_info=True)
            if temp_dir_for_conversion: shutil.rmtree(temp_dir_for_conversion)
            return None, None
        except subprocess.CalledProcessError as e:
            logging.error(f"LibreOffice转换PPT文件 '{ppt_path}' 时出错。命令: '{e.cmd}' 返回码: {e.returncode}", exc_info=True)
            logging.error(f"Stderr: {e.stderr.decode('utf-8', 'ignore') if e.stderr else 'N/A'}")
            logging.error(f"Stdout: {e.stdout.decode('utf-8', 'ignore') if e.stdout else 'N/A'}")
            if temp_dir_for_conversion: shutil.rmtree(temp_dir_for_conversion)
            return None, None
        except Exception as e:
            logging.error(f"转换PPT文件 '{ppt_path}' 时发生未知错误: {str(e)}", exc_info=True)
            if temp_dir_for_conversion: shutil.rmtree(temp_dir_for_conversion)
            return None, None

    def extract_ppt_content(self, file_path_or_url: str) -> Dict:
        """
        提取PPT/PPTX文件的文本内容，支持本地路径和URL。

        :param file_path_or_url: PPT或PPTX文件的本地路径或URL。
        :return: 包含提取内容的字典。
        """
        logging.info(f"开始提取文件内容: {file_path_or_url}")
        content = {
            "文件名": os.path.basename(file_path_or_url), # 初始设为输入，若是URL会被覆盖
            "幻灯片数量": 0,
            "幻灯片内容": [],
            "错误": "" # 用于记录提取阶段的错误信息
        }
        
        # 用于管理临时文件/目录的变量
        processed_file_source: Union[str, io.BytesIO] = file_path_or_url # 可以是路径字符串或BytesIO对象
        temp_dir_for_soffice_output_to_clean = None # soffice转换输出的临时目录
        downloaded_ppt_temp_file_path = None      # 从URL下载的.ppt保存的临时文件路径

        try:
            if self._is_valid_url(file_path_or_url):
                logging.info(f"输入为URL: {file_path_or_url}。正在尝试下载...")
                # 从URL路径中提取原始文件名和扩展名
                parsed_url = urlparse(file_path_or_url)
                original_filename_from_url = os.path.basename(parsed_url.path)
                content["文件名"] = original_filename_from_url # 更新为URL中的文件名
                file_ext_from_url = os.path.splitext(original_filename_from_url)[1].lower()
                logging.info(f"从URL推断文件名: '{original_filename_from_url}', 扩展名: '{file_ext_from_url}'")

                headers = {'User-Agent': 'Mozilla/5.0'} # 简单的User-Agent
                response = requests.get(file_path_or_url, headers=headers, timeout=120) # 2分钟下载超时
                response.raise_for_status() # 检查HTTP错误
                file_content_bytes = response.content
                logging.info(f"成功从URL下载内容，大小: {len(file_content_bytes)}字节。")

                if file_ext_from_url == '.pptx':
                    logging.info("处理从URL下载的 .pptx 内容。")
                    processed_file_source = io.BytesIO(file_content_bytes) # python-pptx可直接处理BytesIO
                elif file_ext_from_url == '.ppt':
                    logging.info("处理从URL下载的 .ppt 内容，将保存到临时文件进行转换。")
                    # 创建一个带特定后缀的命名临时文件，以便soffice能识别
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".ppt") as tmp_ppt_file:
                        tmp_ppt_file.write(file_content_bytes)
                        downloaded_ppt_temp_file_path = tmp_ppt_file.name # 保存路径以备后续清理
                    
                    logging.info(f"从URL下载的 .ppt 内容已保存到临时文件: {downloaded_ppt_temp_file_path}")
                    # 使用这个临时文件的路径进行转换
                    converted_pptx_path, temp_dir_soffice = self.convert_ppt_to_pptx(downloaded_ppt_temp_file_path)
                    if not converted_pptx_path:
                        content["错误"] = "从URL下载的PPT文件转换失败"
                        logging.error(content["错误"] + f" (源URL: {file_path_or_url})")
                        return content # 返回包含错误信息的content字典
                    
                    processed_file_source = converted_pptx_path # 这是转换后的 .pptx 文件路径
                    temp_dir_for_soffice_output_to_clean = temp_dir_soffice # 这个目录是soffice转换输出的地方
                    logging.info(f"URL .ppt ({file_path_or_url}) 已成功转换为 .pptx: {processed_file_source}")
                else:
                    content["错误"] = f"不支持的URL文件扩展名: {file_ext_from_url}"
                    logging.error(content["错误"] + f" (源URL: {file_path_or_url})")
                    return content
            
            else: # 处理本地文件路径
                logging.info(f"输入为本地文件路径: {file_path_or_url}")
                local_file_ext = os.path.splitext(file_path_or_url)[1].lower()
                if local_file_ext == '.ppt':
                    logging.info(f"检测到本地 .ppt 文件，尝试转换为 .pptx: {file_path_or_url}")
                    converted_pptx_path, temp_dir_soffice = self.convert_ppt_to_pptx(file_path_or_url)
                    if not converted_pptx_path:
                        content["错误"] = "本地PPT文件转换失败"
                        logging.error(content["错误"] + f" (文件: {file_path_or_url})")
                        return content
                    processed_file_source = converted_pptx_path
                    temp_dir_for_soffice_output_to_clean = temp_dir_soffice
                    logging.info(f"本地 .ppt 文件已成功转换为 .pptx: {processed_file_source}")
                elif local_file_ext == '.pptx':
                    processed_file_source = file_path_or_url # 直接使用本地pptx路径
                else:
                    content["错误"] = f"不支持的本地文件扩展名: {local_file_ext}"
                    logging.error(content["错误"] + f" (文件: {file_path_or_url})")
                    return content

            # --- 统一使用 python-pptx 解析 processed_file_source ---
            logging.info(f"正在使用 python-pptx 解析: {processed_file_source if isinstance(processed_file_source, str) else '内存中的BytesIO对象'}")
            prs = Presentation(processed_file_source) # Presentation() 可以接受路径字符串或文件类对象
            content["幻灯片数量"] = len(prs.slides)
            logging.info(f"文件共有 {len(prs.slides)} 张幻灯片。")

            for idx, slide in enumerate(prs.slides):
                slide_idx_for_log = idx + 1
                slide_data = {"页码": slide_idx_for_log, "标题": "", "文本内容": [], "备注": ""}
                if slide.shapes.title:
                    slide_data["标题"] = slide.shapes.title.text.strip()
                
                page_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text and shape.text_frame.text.strip():
                        is_title_shape = (slide.shapes.title is not None and shape == slide.shapes.title)
                        if not is_title_shape: # 避免重复添加标题文本
                            page_texts.append(shape.text_frame.text.strip())
                slide_data["文本内容"] = page_texts
                
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame and notes_slide.notes_text_frame.text:
                        slide_data["备注"] = notes_slide.notes_text_frame.text.strip()
                content["幻灯片内容"].append(slide_data)
            
            logging.info(f"内容提取完成: {file_path_or_url}")
            return content

        except requests.exceptions.RequestException as e_req:
            content["错误"] = f"下载URL时出错: {str(e_req)}"
            logging.error(f"下载URL '{file_path_or_url}' 时发生网络错误: {e_req}", exc_info=True)
            return content
        except Exception as e_main:
            content["错误"] = f"提取PPT内容时发生未知错误: {str(e_main)}"
            logging.error(f"提取PPT内容时出错 ('{file_path_or_url}'): {e_main}", exc_info=True)
            return content
        finally:
            # 清理 soffice 转换输出的临时目录 (如果存在)
            if temp_dir_for_soffice_output_to_clean:
                try:
                    shutil.rmtree(temp_dir_for_soffice_output_to_clean)
                    logging.info(f"已清理 soffice 转换输出的临时目录: {temp_dir_for_soffice_output_to_clean}")
                except Exception as e_clean_soffice:
                    logging.error(f"清理 soffice 临时目录 {temp_dir_for_soffice_output_to_clean} 失败: {e_clean_soffice}", exc_info=True)
            
            # 清理从URL下载的临时.ppt文件 (如果存在)
            if downloaded_ppt_temp_file_path and os.path.exists(downloaded_ppt_temp_file_path):
                try:
                    os.remove(downloaded_ppt_temp_file_path)
                    logging.info(f"已清理从URL下载的临时 .ppt 文件: {downloaded_ppt_temp_file_path}")
                except Exception as e_clean_dl:
                    logging.error(f"清理下载的临时 .ppt 文件 {downloaded_ppt_temp_file_path} 失败: {e_clean_dl}", exc_info=True)


    # 示例： summarize_ppt 方法（确保它调用了修改后的 extract_ppt_content）
    def summarize_ppt(self, file_path_or_url: str, custom_prompt: Optional[str] = None) -> str:
        logging.info(f"开始总结PPT: {file_path_or_url}")
        content_data = self.extract_ppt_content(file_path_or_url) # 调用已支持URL的方法

        if content_data.get("错误") or not content_data.get("幻灯片内容"):
            error_msg = content_data.get("错误", "未知提取错误或PPT无有效内容")
            logging.error(f"无法提取或内容为空: {file_path_or_url}。错误: {error_msg}")
            return f"无法处理PPT文件 '{os.path.basename(file_path_or_url)}'。原因: {error_msg}"

        try:
            content_str = json.dumps(content_data, ensure_ascii=False, indent=2)
            logging.info(f"提取的PPT内容字符串长度: {len(content_str)}")
            MAX_CONTENT_STR_LEN = 7500 
            if len(content_str) > MAX_CONTENT_STR_LEN:
                logging.warning(f"PPT内容字符串过长 ({len(content_str)} chars)，将截断发送给LLM。")
                # 简单的截断JSON字符串可能导致其无效，更安全的方式是截断 "幻灯片内容" 里的文本
                # 但为快速实现，这里仅截断整体字符串，并确保JSON结构尽量保持有效
                # 注意：这种截断方式可能导致JSON无效，实际应用中需要更智能的截断策略
                placeholder = "\"...[内容过长，已截断]\"}]}" 
                # 尝试找到一个合适的截断点，比如最后一个幻灯片内容的末尾附近
                # 为了简单，这里还是用之前的截断方式，但提示其风险
                content_str = content_str[:MAX_CONTENT_STR_LEN - len(placeholder)] + placeholder
                logging.debug(f"截断后的内容 (部分): {content_str[:200]}...{content_str[-200:]}")


        except Exception as e:
            logging.error(f"序列化提取的PPT内容时出错: {str(e)}", exc_info=True)
            return "序列化PPT内容时发生内部错误。"

        prompt_to_use = custom_prompt if custom_prompt else """请基于以下提供的PPT演示文稿的提取内容（JSON格式），用中文分析这个演示文稿，请覆盖以下方面：
1.  **核心主题与目标受众**：PPT试图传达的核心信息是什么？主要面向哪些听众？
2.  **结构与逻辑流程**：PPT的组织结构如何？（例如，问题-解决方案，时间顺序，主题分类等）。各部分之间是如何衔接的？
3.  **关键幻灯片内容摘要**：逐页或分章节概括每张（或每组重要）幻灯片的主要观点和信息。
4.  **关键论点/数据/证据**：识别出支撑核心主题的关键论点、重要数据或证据。
5.  **视觉与表达特点**：（如果能从文本推断）PPT在视觉设计或表达上有何特点？（例如，图表为主，文字密集，案例驱动等）。
6.  **潜在的演讲亮点与改进建议**：这份PPT的演讲亮点可能在哪里？有无明显的改进空间（从内容结构、清晰度等方面）？
7.  **总结性概述**：给出一个简短的整体总结。(请确保总结在50字以内)""" # 修改了默认prompt中对字数的要求

        user_content = f"{prompt_to_use}\n\n# 演示文稿提取内容 (JSON):\n```json\n{content_str}\n```"
        # logging.debug(f"发送给 LLM 的 User Content (部分): {user_content[:500]}...") # 内容可能很大

        try:
            logging.info(f"向 LLM ({self.model_name}) 发送请求进行PPT总结...")
            response = self.client.chat.completions.create(
                model=self.model_name,
                messages=[
                    {"role": "system", "content": "你是一位经验丰富的演示文稿分析专家和内容策略师。"},
                    {"role": "user", "content": user_content}
                ],
                temperature=0.3,
                top_p=0.8
            )
            summary = response.choices[0].message.content
            logging.info(f"LLM PPT总结接收成功，长度: {len(summary) if summary else 0}")
            return summary if summary else "未能从模型获取有效的总结内容。"
        except Exception as e:
            logging.error(f"调用 LLM 分析PPT时出错: {str(e)}", exc_info=True)
            return "分析PPT时与AI服务通信失败。"

    def compare_ppts(self, file_paths: List[str], custom_prompt: str = None) -> str:
        """
        使用 LLM 比较多个PPT文件的内容。

        :param file_paths: 包含多个PPT/PPTX文件路径的列表。
        :param custom_prompt: 用户自定义的比较提示词。
        :return: LLM 生成的PPT比较结果文本。
        """
        logging.info(f"开始比较PPT文件: {file_paths}")
        if len(file_paths) < 2:
            logging.warning("需要至少两个PPT文件进行比较。")
            return "错误：需要至少两个PPT文件进行比较。"

        all_contents_data = []
        for file_path in file_paths:
            logging.info(f"正在为比较提取内容: {file_path}")
            content = self.extract_ppt_content(file_path)
            if "错误" not in content and content.get("幻灯片内容"):
                all_contents_data.append(content)
            else:
                logging.warning(f"跳过文件 {file_path}，因提取内容失败或为空。")
        
        if len(all_contents_data) < 2:
            logging.error(f"未能成功提取至少两个PPT的内容进行比较。有效提取数: {len(all_contents_data)}")
            return f"错误：未能成功提取至少两个PPT的内容以供比较（成功提取 {len(all_contents_data)} 个）。"

        try:
            contents_str = json.dumps(all_contents_data, ensure_ascii=False, indent=2)
            logging.info(f"多个PPT提取内容合并后的字符串长度: {len(contents_str)}")
            MAX_CONTENTS_STR_LEN = 7500 # 同样需要注意长度
            if len(contents_str) > MAX_CONTENTS_STR_LEN:
                logging.warning(f"多个PPT内容字符串过长 ({len(contents_str)} chars)，将截断发送给LLM。")
                contents_str = contents_str[:MAX_CONTENTS_STR_LEN] + "\n... [内容已截断]"
        except Exception as e:
            logging.error(f"序列化多个PPT提取内容时出错: {str(e)}", exc_info=True)
            return "序列化多个PPT内容时发生内部错误。"

        prompt_to_use = custom_prompt if custom_prompt else """请基于以下提供的多个PPT演示文稿的提取内容（JSON格式数组），用中文深入比较这些演示文稿的异同点。分析应侧重于：
1.  **核心主题与目标的对比**：各PPT的核心议题是什么？它们的目标受众是否相同或相似？
2.  **内容覆盖与深度的比较**：在共同主题上，各PPT的内容覆盖范围和探讨深度有何不同？有无各自独特的侧重点？
3.  **结构与逻辑流程的差异**：各PPT的组织结构和论证逻辑有何异同？哪种结构更有效？
4.  **关键信息与论据的对比**：它们提出的关键信息、数据或证据是相似还是互补，亦或存在矛盾？
5.  **表达风格与侧重点**：从提取的文本来看，各PPT的表达风格（如正式、非正式、数据驱动、故事化等）有何不同？
6.  **综合评价与建议**：综合来看，这些PPT各自的优势和劣势是什么？如果需要选择一个或整合，你会如何建议？"""
        
        user_content = f"{prompt_to_use}\n\n# 各演示文稿提取内容 (JSON数组):\n```json\n{contents_str}\n```"
        logging.debug(f"发送给 LLM 的 User Content (部分): {user_content[:500]}...")

        try:
            logging.info(f"向 LLM ({self.model_name}) 发送请求进行PPT比较...")
            response = self.client.chat.completions.create(
                model=self.model_name,
                messages=[
                    {"role": "system", "content": "你是一位顶级的演示文稿评审专家，擅长对多个PPT进行深度对比分析和战略评估。"},
                    {"role": "user", "content": user_content}
                ],
                temperature=0.3,
                top_p=0.8
            )
            comparison_summary = response.choices[0].message.content
            logging.info(f"LLM PPT比较结果接收成功，长度: {len(comparison_summary)}")
            return comparison_summary
        except Exception as e:
            logging.error(f"调用 LLM 比较PPT时出错: {str(e)}", exc_info=True)
            return "比较PPT时与AI服务通信失败。"

    def process_directory(self, directory_path: str, output_file: str):
        """
        处理目录中的所有PPT/PPTX文件，并对每个文件生成总结。

        :param directory_path: 包含PPT文件的目录路径。
        :param output_file: 保存所有总结报告的文件路径。
        """
        logging.info(f"开始处理目录: {directory_path}，输出到: {output_file}")
        try:
            dir_path = Path(directory_path)
            if not dir_path.is_dir():
                logging.error(f"提供的路径 '{directory_path}' 不是一个有效的目录。")
                return

            ppt_files = []
            # 使用 rglob 进行递归搜索
            for ext in ['.ppt', '.pptx']:
                ppt_files.extend(list(dir_path.rglob(f'*{ext}')))
            
            if not ppt_files:
                logging.warning(f"在目录 '{directory_path}' 及其子目录中没有找到支持的PPT/PPTX文件。")
                return
            
            logging.info(f"在 '{directory_path}' 中找到 {len(ppt_files)} 个PPT/PPTX文件: {[str(f) for f in ppt_files]}")

            all_summaries_content = []
            for ppt_file_path in tqdm(ppt_files, desc="逐个处理目录中的PPT"):
                logging.info(f"正在处理文件: {str(ppt_file_path)}")
                summary = self.summarize_ppt(str(ppt_file_path))
                # 检查总结是否有效 (不是错误信息)
                if summary and not summary.startswith("无法处理") and not summary.startswith("错误：") and not summary.startswith("分析PPT时与AI服务通信失败"):
                    all_summaries_content.append(f"--- 文件: {ppt_file_path.name} ---\n{summary}\n\n")
                else:
                    all_summaries_content.append(f"--- 文件: {ppt_file_path.name} ---\n总结失败或文件无效。\n详细信息: {summary}\n\n")
                time.sleep(1)  # 控制API调用频率

            if not all_summaries_content:
                logging.info("未能对目录中的任何PPT文件生成有效总结。")
                report_content = "未能对目录中的任何PPT文件生成有效总结。"
            else:
                report_content = "".join(all_summaries_content)

            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(report_content)
            logging.info(f"所有PPT的总结报告已保存到: {output_file}")

        except Exception as e:
            logging.error(f"处理目录 '{directory_path}' 时发生意外错误: {str(e)}", exc_info=True)


# --- 辅助函数：创建示例文件 ---
def create_dummy_pptx(file_path, num_slides=3):
    if not os.path.exists(file_path):
        logging.info(f"创建示例 PPTX 文件: {file_path}")
        prs = Presentation()
        for i in range(num_slides):
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            title.text_frame.text = f"幻灯片 {i+1} 标题"
            
            content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
            p = content_box.text_frame.add_paragraph()
            p.text = f"这是幻灯片 {i+1} 的一些示例文本内容。"
            p.add_run().text = " 重点内容。"
            if i % 2 == 0:
                 p = content_box.text_frame.add_paragraph()
                 p.text = "这里是额外的一行，用于增加多样性。"
            
            if slide.has_notes_slide: # Should always be true for new slides
                 notes_tf = slide.notes_slide.notes_text_frame
                 notes_tf.text = f"这是幻灯片 {i+1} 的演讲者备注。"
        
        # 确保目录存在
        os.makedirs(os.path.dirname(file_path) or '.', exist_ok=True)
        prs.save(file_path)

def create_dummy_ppt(file_path):
    """Creates a dummy empty .ppt file for conversion testing if it doesn't exist."""
    if not os.path.exists(file_path):
        logging.info(f"创建空的占位 .ppt 文件 (用于转换测试): {file_path}")
        os.makedirs(os.path.dirname(file_path) or '.', exist_ok=True)
        with open(file_path, 'w') as f:
            f.write("This is a dummy ppt file for conversion testing.")
        logging.info(f"注意: '{file_path}' 是一个占位符文件，不是一个真正的PPT。LibreOffice应能处理它或报告错误。")

class DocumentSummarizer:
    def __init__(self, api_key: str):
        """
        初始化 DocumentSummarizer。
        :param api_key: 用于访问文心一言 API 的密钥。
        """
        self.client = OpenAI(
            api_key=api_key,
            base_url="https://qianfan.baidubce.com/v2"
        )
        self.chunk_size = 4000
        self.overlap_size = 200
        self.model_name = "ernie-4.5-turbo-128k" # 您指定的模型
        logging.info(
            f"DocumentSummarizer 初始化完成。API base_url: {self.client.base_url}, "
            f"Model: {self.model_name}, Chunk Size: {self.chunk_size}, Overlap Size: {self.overlap_size}"
        )
        logging.warning(
            "注意: 当前 summarize_document 方法会将整个提取的文档内容作为单个JSON字符串发送给LLM。"
            "即使 ernie-4.5-turbo-128k 支持长上下文，对于超大文档，将结构化JSON直接传入仍需注意实际token消耗和成本。"
            "类中定义的 chunk_size, overlap_size 和 split_text 方法目前未用于对提取的文档原始文本内容进行分块总结。"
            "如需更精细或成本优化的超长文档处理，请考虑修改 summarize_document 以实现基于原始文本的分块和可能的分层总结策略。"
        )

    def _is_valid_url(self, url: str) -> bool:
        """简单的URL格式检查，判断是否为 HTTP/HTTPS URL"""
        try:
            result = urlparse(url)
            return all([result.scheme in ['http', 'https'], result.netloc])
        except Exception:
            return False

    def _download_content_from_url(self, url: str) -> Optional[bytes]:
        """从URL下载内容并返回字节。"""
        logging.info(f"尝试从URL下载内容: {url}")
        try:
            headers = {'User-Agent': 'Mozilla/5.0'}
            response = requests.get(url, headers=headers, timeout=120) # 2分钟下载超时
            response.raise_for_status() # 检查HTTP错误
            logging.info(f"成功从URL下载内容，大小: {len(response.content)}字节。")
            return response.content
        except requests.exceptions.RequestException as e_req:
            logging.error(f"下载URL '{url}' 时发生网络错误: {e_req}", exc_info=True)
            return None
        except Exception as e_download:
            logging.error(f"下载或处理URL '{url}' 时发生未知错误: {e_download}", exc_info=True)
            return None


    def extract_pdf_content(self, file_path_or_url: str) -> Dict:
        """
        提取PDF文件的文本内容，支持本地路径和URL。
        :param file_path_or_url: PDF文件的本地路径或URL。
        :return: 包含提取内容的字典。
        """
        logging.info(f"开始提取PDF内容: {file_path_or_url}")
        content_dict = { # 重命名以避免与变量 content 冲突
            "文件名": os.path.basename(file_path_or_url),
            "来源": "本地文件" if not self._is_valid_url(file_path_or_url) else "URL",
            "页数": 0,
            "内容": [],
            "错误": ""
        }
        pdf_stream: Optional[io.BytesIO] = None

        try:
            if self._is_valid_url(file_path_or_url):
                content_dict["文件名"] = os.path.basename(urlparse(file_path_or_url).path) # 从URL获取文件名
                downloaded_bytes = self._download_content_from_url(file_path_or_url)
                if downloaded_bytes is None:
                    content_dict["错误"] = "从URL下载PDF内容失败"
                    return content_dict
                pdf_stream = io.BytesIO(downloaded_bytes)
            else: # 本地文件
                if not os.path.exists(file_path_or_url):
                    content_dict["错误"] = f"本地PDF文件未找到: {file_path_or_url}"
                    logging.error(content_dict["错误"])
                    return content_dict
                pdf_stream = open(file_path_or_url, 'rb') # type: ignore # PyPDF2需要文件对象

            pdf_reader = PyPDF2.PdfReader(pdf_stream) # PyPDF2.PdfReader 可以接受文件对象
            num_pages = len(pdf_reader.pages)
            content_dict["页数"] = num_pages
            logging.info(f"PDF共有 {num_pages} 页。")

            for page_num in range(num_pages):
                try:
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    if text and text.strip():
                        content_dict["内容"].append({
                            "页码": page_num + 1,
                            "文本": text.strip()
                        })
                        logging.debug(f"提取到第 {page_num + 1} 页文本，长度 {len(text.strip())}")
                    else:
                        logging.debug(f"第 {page_num + 1} 页无有效文本内容。")
                except Exception as e_page:
                    logging.warning(f"提取PDF第 {page_num + 1} 页内容时出错: {str(e_page)}", exc_info=False) # Debug时可设为True
                    content_dict["内容"].append({
                        "页码": page_num + 1,
                        "文本": f"[错误：无法提取此页内容 - {str(e_page)}]"
                    })
            
            if not content_dict["内容"] and num_pages > 0 : # 有页面但是没提取到内容
                 logging.warning(f"未能从PDF '{file_path_or_url}' 中提取任何文本内容（可能是图片型PDF或加密PDF）。")
                 content_dict["错误"] = "未能提取任何文本内容（可能是图片型PDF或加密PDF）"
            elif num_pages == 0:
                 logging.warning(f"PDF '{file_path_or_url}' 不包含任何页面。")
                 content_dict["错误"] = "PDF文件不包含任何页面"
            else:
                 logging.info(f"PDF内容提取完成: {file_path_or_url}")
            return content_dict

        except PyPDF2.errors.PdfReadError as e_pdf_read: # 更具体的PyPDF2错误
            logging.error(f"读取PDF '{file_path_or_url}' 时出错 (可能已损坏或加密): {e_pdf_read}", exc_info=True)
            content_dict["错误"] = f"读取PDF文件失败 (可能已损坏或加密): {e_pdf_read}"
            return content_dict
        except Exception as e:
            logging.error(f"提取PDF '{file_path_or_url}' 内容时发生严重错误: {e}", exc_info=True)
            content_dict["错误"] = f"提取PDF内容时出错: {e}"
            return content_dict
        finally:
            if pdf_stream and hasattr(pdf_stream, 'close'): # 如果是从文件打开的，需要关闭
                pdf_stream.close()


    def extract_docx_content(self, file_path_or_url: str) -> Dict:
        """
        提取DOCX文件的文本内容，按标题分段，支持本地路径和URL。
        :param file_path_or_url: DOCX文件的本地路径或URL。
        :return: 包含提取内容的字典。
        """
        logging.info(f"开始提取DOCX内容: {file_path_or_url}")
        content_dict = {
            "文件名": os.path.basename(file_path_or_url),
            "来源": "本地文件" if not self._is_valid_url(file_path_or_url) else "URL",
            "总段落数": 0,
            "内容": [],
            "错误": ""
        }
        docx_source: Union[str, io.BytesIO] = file_path_or_url

        try:
            if self._is_valid_url(file_path_or_url):
                content_dict["文件名"] = os.path.basename(urlparse(file_path_or_url).path)
                downloaded_bytes = self._download_content_from_url(file_path_or_url)
                if downloaded_bytes is None:
                    content_dict["错误"] = "从URL下载DOCX内容失败"
                    return content_dict
                docx_source = io.BytesIO(downloaded_bytes)
            else: # 本地文件
                if not os.path.exists(file_path_or_url):
                    content_dict["错误"] = f"本地DOCX文件未找到: {file_path_or_url}"
                    logging.error(content_dict["错误"])
                    return content_dict
                docx_source = file_path_or_url # Document()可以直接处理路径

            doc = DocxDocument(docx_source) # DocxDocument可以接受路径字符串或文件类对象
            content_dict["总段落数"] = len(doc.paragraphs)
            logging.info(f"DOCX共有 {len(doc.paragraphs)} 个原始段落。")

            current_section = {"标题": "默认起始段落", "文本": []}
            first_para_processed = False
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                is_heading = para.style.name.lower().startswith('heading')
                
                if is_heading:
                    if current_section["文本"] or (current_section["标题"] != "默认起始段落" and not content_dict["内容"]): # 如果有文本，或者这是第一个自定义标题（即使没文本）
                        content_dict["内容"].append(current_section)
                        logging.debug(f"保存DOCX章节: '{current_section['标题']}', 文本段落数: {len(current_section['文本'])}")
                    current_section = {"标题": text, "文本": []}
                else:
                    current_section["文本"].append(text)
                first_para_processed = True
            
            if current_section["文本"] or (not content_dict["内容"] and current_section["标题"] != "默认起始段落"):
                content_dict["内容"].append(current_section)
                logging.debug(f"保存最后一个DOCX章节: '{current_section['标题']}', 文本段落数: {len(current_section['文本'])}")

            if not content_dict["内容"]:
                 logging.warning(f"未能从DOCX '{file_path_or_url}' 中提取任何结构化文本内容。")
                 all_text = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
                 if all_text:
                     content_dict["内容"].append({"标题": "所有文本（无明确结构）", "文本": all_text})
                     logging.info("DOCX无明确结构或标准标题，已提取所有文本段落。")
                 else:
                    content_dict["错误"] = "未能提取任何文本内容"
            
            logging.info(f"DOCX内容提取完成: {file_path_or_url}")
            return content_dict

        except Exception as e:
            logging.error(f"提取DOCX '{file_path_or_url}' 内容时出错: {e}", exc_info=True)
            content_dict["错误"] = f"提取DOCX内容时出错: {e}"
            return content_dict

    # split_text, summarize_document, compare_documents, process_directory 方法与您之前提供的一致
    # 它们会调用更新后的 extract_pdf_content 和 extract_docx_content
    # 以下是 summarize_document 的一个片段，展示如何调用
    
    def split_text(self, text: str) -> List[str]:
        """
        将长文本分割成重叠的块。
        修正了start指针更新逻辑以确保前进。

        :param text: 待分割的文本。
        :return: 分割后的文本块列表。
        """
        if not text:
            logging.warning("输入文本为空，无法执行 split_text。")
            return []
            
        logging.info(f"开始分割文本，总长度: {len(text)}，Chunk Size: {self.chunk_size}, Overlap: {self.overlap_size}")
        chunks = []
        start = 0
        text_length = len(text)

        while start < text_length:
            end = min(start + self.chunk_size, text_length)
            actual_chunk_end = end

            if end < text_length:
                search_start_offset = min(start + self.overlap_size // 2, end -1) # 调整搜索起点
                last_period = text.rfind('。', search_start_offset, end)
                last_newline = text.rfind('\n', search_start_offset, end)
                
                split_point = -1
                if last_period != -1: split_point = max(split_point, last_period)
                if last_newline != -1: split_point = max(split_point, last_newline)

                if split_point > start: 
                    actual_chunk_end = split_point + 1
                    logging.debug(f"找到分割点 {split_point}，块 [{start}:{actual_chunk_end}]")
                else:
                    logging.debug(f"未在 [{search_start_offset}:{end}] 找到合适的自然分割点，使用默认块结束位置 {actual_chunk_end}")
            
            chunks.append(text[start:actual_chunk_end])
            logging.debug(f"添加块: text[{start}:{actual_chunk_end}], 长度 {len(chunks[-1])}")

            if actual_chunk_end == text_length:
                logging.info("已到达文本末尾，split_text 结束。")
                break
            
            next_start_candidate = actual_chunk_end - self.overlap_size
            
            if next_start_candidate <= start:
                logging.warning(
                    f"split_text: 下一个起始点 ({next_start_candidate}) <= 当前起始点 ({start}). "
                    f"块长度 {actual_chunk_end - start}, 重叠 {self.overlap_size}. "
                    f"将起始点强制前进到 {actual_chunk_end} (牺牲此步重叠) 以避免循环。"
                )
                start = actual_chunk_end
            else:
                start = next_start_candidate
            logging.debug(f"下一个块的起始位置: {start}")
        
        logging.info(f"文本分割完成，共 {len(chunks)} 个块。")
        return chunks

    def summarize_document(self, file_path_or_url: str, custom_prompt: Optional[str] = None) -> str:
        logging.info(f"开始总结文档: {file_path_or_url}")
        file_ext = ""
        if self._is_valid_url(file_path_or_url):
            # 从URL路径中提取扩展名
            parsed_url_path = urlparse(file_path_or_url).path
            file_ext = os.path.splitext(parsed_url_path)[1].lower()
        else:
            file_ext = os.path.splitext(file_path_or_url)[1].lower()
        
        content_data: Dict = {}

        if file_ext == '.pdf':
            content_data = self.extract_pdf_content(file_path_or_url)
        elif file_ext == '.docx':
            content_data = self.extract_docx_content(file_path_or_url)
        else:
            logging.warning(f"不支持的文件格式或无法从URL推断格式: '{file_ext}' (源: {file_path_or_url})")
            return f"不支持的文件格式或无法从URL推断格式: '{file_ext}'"

        if content_data.get("错误") or not content_data.get("内容"):
            error_msg = content_data.get("错误", "未知提取错误或文档无内容")
            logging.error(f"无法提取或内容为空: {file_path_or_url}。错误: {error_msg}")
            return f"无法处理文档 '{content_data.get('文件名', os.path.basename(file_path_or_url))}'。原因: {error_msg}"

        try:
            content_str = json.dumps(content_data, ensure_ascii=False, indent=2)
            logging.info(f"提取的文档内容序列化后字符串长度: {len(content_str)}")
            
            # 即使是128k模型，也应有实际的token限制。JSON结构会比纯文本消耗更多token。
            # 这里的7800字符是一个非常粗略的估计，实际应根据模型文档和测试调整。
            # 对于 ernie-4.5-turbo-128k，理论上可以处理更长的内容，但仍需注意API的实际限制和成本。
            # 这里的MAX_CONTENT_STR_LEN主要是为了防止意外的超大请求。
            MAX_CONTENT_STR_LEN = 100000 # 增大这个值，因为模型是128k tokens
            if len(content_str) > MAX_CONTENT_STR_LEN:
                logging.warning(
                    f"文档内容JSON字符串过长 ({len(content_str)} chars)，将截断至约 {MAX_CONTENT_STR_LEN} 字符后发送给LLM。"
                    "对于 ernie-4.5-turbo-128k，这仍可能在token限制内，但请监控实际消耗。"
                    "如需最优化，请考虑基于原始文本的分块总结。"
                )
                # 截断JSON字符串时要小心，确保其至少是部分有效的或模型能容忍的。
                # 一个更安全的截断方式是减少 "内容" 数组中的元素或截断元素内的 "文本"。
                # 为简单起见，仍进行字符串截断，但对于128k模型，这可能不是主要瓶颈。
                placeholder_text = "\"...[内容JSON因过长被截断，但仍保留了大部分结构]\"}]}" # 确保json结构在截断后仍有意义
                cutoff_point = MAX_CONTENT_STR_LEN - len(placeholder_text)
                if cutoff_point > 0 :
                    content_str = content_str[:cutoff_point] + placeholder_text
                else: # 如果最大长度甚至容不下占位符，就直接截断
                    content_str = content_str[:MAX_CONTENT_STR_LEN]
        except Exception as e:
            logging.error(f"序列化提取的文档内容时出错: {e}", exc_info=True)
            return "序列化文档内容时发生内部错误。"

        prompt_to_use = custom_prompt if custom_prompt else """请基于以下提供的文档提取内容（JSON格式），用中文全面分析此文档，包括但不限于：
1.  **核心议题与目的**：文档主要讨论什么？其核心目标或意图是什么？
2.  **结构与组织**：文档的整体结构是怎样的？（例如，引言、主体章节、结论；问题-分析-方案等）。
3.  **主要章节/部分摘要**：针对提取内容中的各主要部分（如PDF的页面、DOCX的章节标题下的文本），分别概括其要点。
4.  **关键信息、论点与证据**：识别文档中提出的关键信息、核心论点以及支撑这些论点的主要数据或证据。
5.  **目标读者与潜在影响**：这份文档可能面向哪些读者？它可能产生哪些影响或达到什么效果？
6.  **总结性概述**：对整个文档给出一个精炼的总体总结。(请确保总结在50字以内)""" # 调整了默认prompt对总结字数的要求

        user_content = f"{prompt_to_use}\n\n# 文档提取内容 (JSON):\n```json\n{content_str}\n```"

        try:
            logging.info(f"向 LLM ({self.model_name}) 发送请求进行文档总结...")
            response = self.client.chat.completions.create(
                model=self.model_name,
                messages=[
                    {"role": "system", "content": "你是一位资深的文档分析和信息提炼专家，能够精确把握文档核心并生成高质量摘要。"},
                    {"role": "user", "content": user_content}
                ],
                temperature=0.3,
                top_p=0.8
                # stream=True # 如果需要流式输出，可以启用
            )
            summary = ""
            # if stream:
            # for chunk in response:
            # if chunk.choices[0].delta.content:
            # summary_part = chunk.choices[0].delta.content
            # print(summary_part, end="", flush=True) # 实时打印
            # summary += summary_part
            # else:
            summary = response.choices[0].message.content if response.choices[0].message else ""
            
            logging.info(f"LLM 文档总结接收成功，长度: {len(summary) if summary else 0}")
            return summary if summary else "未能从模型获取有效的总结内容。"
        except Exception as e:
            logging.error(f"调用 LLM ({self.model_name}) 分析文档时出错: {e}", exc_info=True)
            return f"分析文档时与AI服务通信失败: {str(e)}"

    # compare_documents 和 process_directory 方法与您之前提供的一致
    # 它们应该能够正确地调用更新后的 summarize_document，进而处理URL
    # (此处省略，因为它们逻辑上依赖 summarize_document 的输入处理能力)
    def compare_documents(self, file_paths_or_urls: List[str], custom_prompt: Optional[str] = None) -> str:
        logging.info(f"开始比较多个文档: {file_paths_or_urls}")
        if len(file_paths_or_urls) < 2:
            logging.warning("需要至少两个文档文件进行比较。")
            return "错误：需要至少两个文档文件进行比较。"

        all_contents_data = []
        for item_path in file_paths_or_urls:
            logging.info(f"正在为比较提取内容: {item_path}")
            file_ext = ""
            if self._is_valid_url(item_path):
                parsed_url_path = urlparse(item_path).path
                file_ext = os.path.splitext(parsed_url_path)[1].lower()
            else:
                file_ext = os.path.splitext(item_path)[1].lower()
            
            content_data = {}
            if file_ext == '.pdf':
                content_data = self.extract_pdf_content(item_path)
            elif file_ext == '.docx':
                content_data = self.extract_docx_content(item_path)
            else:
                logging.warning(f"跳过不支持的文件格式或无法从URL推断格式进行比较: {item_path}")
                continue
            
            if "错误" not in content_data and content_data.get("内容"):
                all_contents_data.append(content_data)
            else:
                logging.warning(f"跳过文件 {item_path}，因提取内容失败或为空。错误详情: {content_data.get('错误')}")
        
        if len(all_contents_data) < 2:
            logging.error(f"未能成功提取至少两个文档的内容进行比较。有效提取数: {len(all_contents_data)}")
            return f"错误：未能成功提取至少两个文档的内容以供比较（成功提取 {len(all_contents_data)} 个）。"

        try:
            contents_str = json.dumps(all_contents_data, ensure_ascii=False, indent=2)
            logging.info(f"多个文档提取内容合并后的字符串长度: {len(contents_str)}")
            MAX_CONTENTS_STR_LEN = 100000 # 根据128k模型调整
            if len(contents_str) > MAX_CONTENTS_STR_LEN:
                logging.warning(f"多个文档内容字符串过长 ({len(contents_str)} chars)，将截断。可能影响比较质量。")
                placeholder_text = "\"...[内容JSON数组因过长被截断]...\"}]"
                cutoff_point = MAX_CONTENTS_STR_LEN - len(placeholder_text)
                if cutoff_point > 0:
                    contents_str = contents_str[:cutoff_point] + placeholder_text
                else:
                    contents_str = contents_str[:MAX_CONTENTS_STR_LEN]
        except Exception as e:
            logging.error(f"序列化多个文档提取内容时出错: {e}", exc_info=True)
            return "序列化多个文档内容时发生内部错误。"

        prompt_to_use = custom_prompt if custom_prompt else """请基于以下提供的多个文档的提取内容（JSON格式数组），用中文深入比较这些文档的异同点。分析应包括：
1.  **核心主题与目标的对比**：各文档的核心议题和主要目标分别是什么？有何异同？
2.  **内容覆盖范围与深度的比较**：在共同或相关主题上，各文档的内容覆盖广度和探讨深度如何？各自有无独特的侧重点或遗漏？
3.  **结构与论证逻辑的差异**：各文档的组织结构（如章节安排、论点呈现顺序）和逻辑推演方式有何不同？
4.  **关键信息、论点与证据的对比**：它们提出的关键信息、核心论点及支撑证据是相似、互补还是存在差异甚至矛盾？
5.  **受众与风格的判断**：从内容和表达推断，各文档可能的目标受众和沟通风格（如学术、商业、技术、通俗等）有何不同？
6.  **综合评价与洞察**：综合来看，这些文档在信息传递、论证力度、完整性等方面各自的优势和不足是什么？能从中获得哪些交叉验证的结论或新的洞察？"""
        
        user_content = f"{prompt_to_use}\n\n# 各文档提取内容 (JSON数组):\n```json\n{contents_str}\n```"

        try:
            logging.info(f"向 LLM ({self.model_name}) 发送请求进行文档比较...")
            response = self.client.chat.completions.create(
                model=self.model_name,
                messages=[
                    {"role": "system", "content": "你是一位顶级的跨文档分析专家，擅长从多个复杂文档中找出深层联系与差异。"},
                    {"role": "user", "content": user_content}
                ],
                temperature=0.3,
                top_p=0.8
            )
            comparison_summary = response.choices[0].message.content if response.choices[0].message else ""
            logging.info(f"LLM 文档比较结果接收成功，长度: {len(comparison_summary)}")
            return comparison_summary if comparison_summary else "未能从模型获取有效的比较结果。"
        except Exception as e:
            logging.error(f"调用 LLM ({self.model_name}) 比较文档时出错: {e}", exc_info=True)
            return f"比较文档时与AI服务通信失败: {str(e)}"

    def process_directory(self, directory_path: str, output_file: str):
        logging.info(f"开始处理目录: {directory_path}，输出到: {output_file}")
        try:
            dir_path = Path(directory_path)
            if not dir_path.is_dir():
                logging.error(f"提供的路径 '{directory_path}' 不是一个有效的目录。")
                return

            doc_files = []
            for ext in ['.pdf', '.docx']:
                doc_files.extend(list(dir_path.rglob(f'*{ext}')))
            
            if not doc_files:
                logging.warning(f"在目录 '{directory_path}' 及其子目录中没有找到支持的文档文件 (.pdf, .docx)。")
                Path(output_file).write_text(f"目录 '{directory_path}' 中未找到支持的文档文件。\n", encoding='utf-8')
                return
            
            logging.info(f"在 '{directory_path}' 中找到 {len(doc_files)} 个文档文件: {[str(f) for f in doc_files]}")

            all_summaries_content = []
            successful_summaries = 0
            for doc_file_path_obj in tqdm(doc_files, desc="逐个处理目录中的文档"):
                doc_file_path = str(doc_file_path_obj) # 转为字符串
                logging.info(f"正在处理文件: {doc_file_path}")
                summary = self.summarize_document(doc_file_path)
                
                # 更新错误检查逻辑
                is_error = False
                error_keywords = ["不支持的文件格式", "无法处理文档", "错误：", "分析文档时与AI服务通信失败", "序列化文档内容时发生内部错误", "未能从模型获取有效的总结内容"]
                if not summary: # 空总结也是一种失败
                    is_error = True
                else:
                    for keyword in error_keywords:
                        if keyword in summary:
                            is_error = True
                            break
                
                if not is_error:
                    all_summaries_content.append(f"--- 文件: {doc_file_path_obj.name} ---\n{summary}\n\n")
                    successful_summaries +=1
                else:
                    all_summaries_content.append(f"--- 文件: {doc_file_path_obj.name} ---\n总结失败或文件无效/不支持。\n详细信息: {summary or '无返回内容'}\n\n")
                
                if len(doc_files) > 1 and doc_file_path_obj != doc_files[-1]:
                     time.sleep(1)

            report_header = f"文档批量处理总结报告\n处理文档总数: {len(doc_files)}\n成功总结数量: {successful_summaries}\n\n"
            final_report_content = report_header + "".join(all_summaries_content)

            Path(output_file).write_text(final_report_content, encoding='utf-8')
            logging.info(f"所有文档的总结报告已保存到: {output_file}")

        except Exception as e:
            logging.error(f"处理目录 '{directory_path}' 时发生意外错误: {e}", exc_info=True)

class VideoSummarizer:
    def __init__(self, api_key: str, base_url: str):
        """
        初始化 VideoSummarizer。
        :param api_key: 用于访问 LLM API 的密钥。
        :param base_url: LLM API 的基础 URL。
        """
        self.client = OpenAI(
            api_key=api_key,
            base_url=base_url
        )
        self.model_name = "qwen-vl-max-latest" # 您指定的模型
        logging.info(f"VideoSummarizer 初始化完成。API base_url: {base_url}, Model: {self.model_name}")
        logging.warning("当前版本主要支持通过 URL 处理视频。本地视频文件处理功能尚未完全实现。")

    def is_valid_url(self, url: str) -> bool:
        """检查URL是否有效"""
        logging.debug(f"检查URL有效性: {url}")
        try:
            result = urlparse(url)
            is_valid = all([result.scheme, result.netloc])
            logging.debug(f"URL '{url}' 解析结果: scheme='{result.scheme}', netloc='{result.netloc}'. 有效性: {is_valid}")
            return is_valid
        except (ValueError, AttributeError) as e: # 更具体的异常捕获
            logging.warning(f"解析URL '{url}' 时发生错误: {e}")
            return False

    def is_valid_file(self, file_path: str) -> bool:
        """检查文件是否存在且可能是视频文件（基于扩展名）"""
        logging.debug(f"检查文件路径有效性: {file_path}")
        if not os.path.exists(file_path):
            logging.warning(f"文件路径不存在: {file_path}")
            return False
        if not os.path.isfile(file_path): # 确保是文件而不是目录
            logging.warning(f"路径不是一个文件: {file_path}")
            return False
            
        video_extensions = ['.mp4', '.avi', '.mov', '.wmv', '.flv', '.mkv', '.mpeg', '.mpg', '.webm']
        ext = os.path.splitext(file_path)[1].lower()
        is_video = ext in video_extensions
        logging.debug(f"文件 '{file_path}' 扩展名: '{ext}'. 是否为视频扩展名: {is_video}")
        return is_video

    def summarize_video(self, video_source: str, prompt: Optional[str] = None) -> str:
        """
        总结单个视频内容。
        
        Args:
            video_source: 视频URL或本地文件路径。
            prompt: 用户自定义的总结提示词。
        Returns:
            视频的总结文本，或错误信息。
        """
        logging.info(f"开始处理单个视频总结: {video_source}")
        video_url_for_api: Optional[str] = None

        if self.is_valid_url(video_source):
            video_url_for_api = video_source
            logging.info(f"输入源为有效URL: {video_url_for_api}")
        elif self.is_valid_file(video_source):
            logging.warning(f"检测到本地文件路径: {video_source}。当前版本暂不支持直接处理本地视频文件上传。")
            # TODO: 此处需要实现将本地文件上传到可被API访问的URL的逻辑,
            #       或者如果API支持直接上传二进制文件流，则调整API调用方式。
            #       目前, 我们将简单返回一个提示信息。
            # video_url_for_api = self._upload_local_file(video_source) # 假设有这个方法
            # if not video_url_for_api:
            #     return "本地视频文件处理失败（上传或转换URL失败）。"
            return "⚠️ 暂不支持本地视频文件直接处理，请提供可公开访问的视频URL。"
        else:
            logging.warning(f"提供的视频源无效: {video_source}")
            return "❌ 无效的视频源，请提供有效的URL或本地文件路径。"

        if prompt is None:
            prompt = """请基于提供的视频，用中文详细分析其内容，包括但不限于：
1.  视频的核心主题和主要讲述的内容是什么？
2.  视频中包含哪些关键的场景、事件或转折点？
3.  是否有重要的人物出现？他们有何显著的动作、言语或表情？
4.  视频的整体氛围（例如：轻松、紧张、悲伤、鼓舞人心等）和视觉风格（例如：纪录片、动画、电影感、Vlog等）是怎样的？
5.  视频可能的目标受众是谁？它试图传递什么信息或达到什么目的？

要求：
1. 输出内容必须控制在50字以内
2. 只提取最核心的信息
3. 使用简洁的语言表达
"""
            logging.debug("使用默认提示词进行视频总结。")
        else:
            logging.debug(f"使用自定义提示词进行视频总结: '{prompt[:100]}...'")

        try:
            logging.info(f"向模型 '{self.model_name}' 发送视频总结请求，视频URL: {video_url_for_api}")
            messages_payload = [
                {
                    "role": "system",
                    "content": [{"type": "text", "text": "你是一位专业的视频内容分析师，能够深入理解视频的视觉和叙事元素，并提供富有洞察力的总结。"}]
                },
                {
                    "role": "user",
                    "content": [
                        {"type": "video_url", "video_url": {"url": video_url_for_api}},
                        {"type": "text", "text": prompt}
                    ]
                }
            ]
            
            response = self.client.chat.completions.create(
                model=self.model_name,
                messages=messages_payload # type: ignore
            )
            
            summary = response.choices[0].message.content
            logging.info(f"成功从模型获取视频总结。总结长度: {len(summary) if summary else 0} 字符。")
            return summary if summary else "未能从模型获取有效的总结内容。"
            
        except Exception as e:
            logging.error(f"使用模型分析视频 '{video_url_for_api}' 时出错: {e}", exc_info=True)
            return f"分析视频时发生错误: {str(e)}"

    def compare_videos(self, video_sources: List[str], prompt: Optional[str] = None) -> str:
        """
        比较多段视频内容。
        
        Args:
            video_sources: 包含多个视频URL或本地文件路径的列表。
            prompt: 用户自定义的比较提示词。
        Returns:
            视频的比较结果文本，或错误信息。
        """
        logging.info(f"开始处理视频比较，共 {len(video_sources)} 个源。")
        if len(video_sources) < 2:
            logging.warning("需要至少两个视频进行比较，当前提供数量不足。")
            return "错误：需要至少两个视频进行比较。"
            
        valid_video_urls_for_api: List[Dict[str, Dict[str, str]]] = [] # 存储 {"type": "video_url", "video_url": {"url": ...}}
        
        for source_idx, source in enumerate(video_sources):
            logging.debug(f"处理比较列表中的源 #{source_idx + 1}: {source}")
            if self.is_valid_url(source):
                valid_video_urls_for_api.append({"type": "video_url", "video_url": {"url": source}})
                logging.info(f"有效URL已添加到比较列表: {source}")
            elif self.is_valid_file(source):
                logging.warning(f"比较功能暂不支持本地视频文件: {source}。此视频将从比较中排除。")
                # TODO: 实现本地文件上传并获取URL的逻辑
            else:
                logging.warning(f"无效的视频源，已从比较列表中排除: {source}")
        
        if len(valid_video_urls_for_api) < 2:
            logging.warning(f"经过验证，有效视频URL数量 ({len(valid_video_urls_for_api)}) 不足两个，无法进行比较。")
            return f"错误：没有足够的有效视频URL（至少需要2个，实际有效 {len(valid_video_urls_for_api)} 个）进行比较。"
            
        if prompt is None:
            prompt = """请基于提供的多段视频，用中文详细对比分析它们之间的异同点，包括但不限于：
1.  **内容与主题**：各视频的核心主题是什么？它们在叙事内容上有何相似之处和明显差异？
2.  **关键场景与事件**：对比各视频中的关键场景或重要事件，它们是如何呈现的？有无对应或对比关系？
3.  **人物与表现**：如果视频中有人物，对比他们在不同视频中的角色、行为或情感表达。
4.  **视觉风格与氛围**：各视频的拍摄手法、色调、剪辑节奏、背景音乐等视觉和听觉元素有何不同？营造的整体氛围有何差异？
5.  **目标与信息传递**：这些视频各自试图传递什么信息或达到什么目的？它们在实现这些目标的方式上有何不同？"""
            logging.debug("使用默认提示词进行视频比较。")
        else:
            logging.debug(f"使用自定义提示词进行视频比较: '{prompt[:100]}...'")
            
        try:
            # 构建单个user消息，包含所有视频URL和文本提示
            user_content_list: List[Dict[str, any]] = []
            user_content_list.extend(valid_video_urls_for_api) # 添加所有视频URL对象
            user_content_list.append({"type": "text", "text": prompt}) # 添加文本提示

            messages_payload = [
                {
                    "role": "system",
                    "content": [{"type": "text", "text": "你是一位顶级的视频对比分析专家，能够敏锐地捕捉多个视频间的细微差别和深层联系。"}]
                },
                {
                    "role": "user",
                    "content": user_content_list
                }
            ]
            
            logging.info(f"向模型 '{self.model_name}' 发送 {len(valid_video_urls_for_api)} 个视频的比较请求。")
            response = self.client.chat.completions.create(
                model=self.model_name,
                messages=messages_payload # type: ignore
            )
            
            comparison_summary = response.choices[0].message.content
            logging.info(f"成功从模型获取视频比较结果。结果长度: {len(comparison_summary) if comparison_summary else 0} 字符。")
            return comparison_summary if comparison_summary else "未能从模型获取有效的比较结果。"
            
        except Exception as e:
            logging.error(f"使用模型比较视频时出错: {e}", exc_info=True)
            return f"比较视频时发生错误: {str(e)}"

    def process_video_list(self, video_sources: List[str], output_file: str):
        """
        处理视频URL列表，对每个视频生成总结，并将所有总结写入一个报告文件。
        
        Args:
            video_sources: 包含多个视频URL或本地文件路径的列表。
            output_file: 保存所有总结报告的文件路径。
        """
        logging.info(f"开始批量处理视频列表，共 {len(video_sources)} 个源，输出到: {output_file}")
        if not video_sources:
            logging.warning("输入的视频源列表为空，无需处理。")
            Path(output_file).write_text("输入的视频源列表为空。\n", encoding='utf-8')
            return

        all_summaries_content: List[str] = []
        successful_summaries = 0
        
        for idx, video_source in enumerate(tqdm(video_sources, desc="批量处理视频中")):
            logging.info(f"正在处理列表中的视频 #{idx + 1}: {video_source}")
            summary = self.summarize_video(video_source) # summarize_video内部已处理URL/文件路径判断
            
            # 为报告确定源的名称
            source_name = ""
            if self.is_valid_url(video_source):
                source_name = video_source
            elif self.is_valid_file(video_source): # 尽管目前不支持处理，但仍记录文件名
                source_name = os.path.basename(video_source)
            else:
                source_name = f"无效源 ({video_source})"

            if summary and not summary.startswith("❌") and not summary.startswith("⚠️") and not summary.startswith("分析视频时发生错误"):
                all_summaries_content.append(f"--- 视频源: {source_name} ---\n{summary}\n\n")
                successful_summaries +=1
            else:
                all_summaries_content.append(f"--- 视频源: {source_name} ---\n总结失败或跳过。\n详细信息: {summary}\n\n")
            
            if idx < len(video_sources) - 1: # 如果不是最后一个视频，则等待
                logging.debug(f"等待1秒以控制API调用频率...")
                time.sleep(1) 

        report_header = f"视频批量处理总结报告\n处理视频总数: {len(video_sources)}\n成功总结数量: {successful_summaries}\n\n"
        final_report_content = report_header + "".join(all_summaries_content)

        try:
            Path(output_file).write_text(final_report_content, encoding='utf-8')
            logging.info(f"视频批量处理完成。总结报告已保存到: {output_file}")
        except Exception as e:
            logging.error(f"保存总结报告到 '{output_file}' 时出错: {e}", exc_info=True)



if __name__ == "__main__":
    # main()
    pass