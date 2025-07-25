import pdfplumber
import docx
import pandas as pd
from pptx import Presentation
import subprocess
import os
import threading
import queue
import time

def convert_to_modern_format(input_path, target_ext):
    """
    用libreoffice将老格式文件（.doc/.ppt）转换为新格式（.docx/.pptx）
    返回转换后的文件路径
    """
    output_dir = os.path.dirname(input_path)
    cmd = [
        "libreoffice",
        "--headless",
        "--convert-to", target_ext,
        "--outdir", output_dir,
        input_path
    ]
    subprocess.run(cmd, check=True)
    base = os.path.splitext(os.path.basename(input_path))[0]
    new_file = os.path.join(output_dir, base + "." + target_ext)
    return new_file

def read_pdf(file_path):
    with pdfplumber.open(file_path) as pdf:
        return "\n".join(page.extract_text() or "" for page in pdf.pages)

def read_word(file_path):
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def read_excel(file_path):
    df = pd.read_excel(file_path, sheet_name=None)
    content = []
    for sheet, data in df.items():
        content.append(f"【Sheet: {sheet}】\n{data.to_string(index=False)}")
    return "\n\n".join(content)

def read_ppt(file_path):
    prs = Presentation(file_path)
    content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content.append(shape.text)
    return "\n".join(content)

def read_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

# 全局队列和线程
convert_queue = queue.Queue()
convert_result = {}
convert_lock = threading.Lock()

def convert_worker():
    while True:
        task_id, input_path, target_ext = convert_queue.get()
        try:
            new_file = convert_to_modern_format(input_path, target_ext)
            with convert_lock:
                convert_result[task_id] = (True, new_file)
        except Exception as e:
            with convert_lock:
                convert_result[task_id] = (False, str(e))
        finally:
            convert_queue.task_done()

# 启动后台线程
threading.Thread(target=convert_worker, daemon=True).start()

def convert_with_queue(input_path, target_ext, timeout=60):
    """
    将转换任务加入队列，等待串行处理，返回转换后文件路径
    """
    task_id = f"{input_path}-{time.time()}"
    with convert_lock:
        convert_result.pop(task_id, None)
    convert_queue.put((task_id, input_path, target_ext))
    start = time.time()
    while True:
        with convert_lock:
            if task_id in convert_result:
                success, result = convert_result.pop(task_id)
                if success:
                    return result
                else:
                    raise Exception(result)
        if time.time() - start > timeout:
            raise TimeoutError("libreoffice转换超时")
        time.sleep(0.2)